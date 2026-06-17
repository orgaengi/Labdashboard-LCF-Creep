"""
generators.py — All generate functions, no tkinter dependency
"""
import os, sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import lab_core as core


# ============================================================
# generate_lcf_creep
# ============================================================

"""
Tool 1 — LCF & Creep Lab Dashboard
====================================
Tracks LCF (capacity=50/yr) and Creep (capacity=22/yr) labs.
Run: python tool1_lcf_creep.py
"""

import os, sys, threading


import openpyxl

# Add current directory to path so lab_core is found



# ── Tool-specific config ──────────────────────────────────
TOOL_TITLE   = "LCF & Creep Lab Dashboard"
LCF_ALLOWED_TYPES = ["LCF", "Creep"]          # unique name: avoids collision with coating
LCF_OUTPUT_FILE = "LCF_Creep_Dashboard.xlsx"   # unique name — module-level
                                                # OUTPUT_FILE is redefined by
                                                # Tool 2 and Tool 4 below, which
                                                # would otherwise silently win
                                                # at call time (same class of
                                                # bug as the v13 ALLOWED_TYPES fix)

DEFAULT_CAPS = {
    "LCF":   50,   # 50 samples/year max
    "Creep": 22,   # 22 samples/year max
}

# ═══════════════════════════════════════════════════════════
#  DASHBOARD GENERATOR
# ═══════════════════════════════════════════════════════════

def generate_lcf_creep(input_path, capacities, theme, progress_cb=None):
    def p(msg):
        if progress_cb: progress_cb(msg)

    p("📖 Reading file…")
    paths = input_path if isinstance(input_path, (list, tuple)) else [input_path]
    if len(paths) > 1:
        df, col_map, errors, warns = merge_files(paths, allowed_types=LCF_ALLOWED_TYPES)
    else:
        df, col_map, errors, warns = core.load_and_filter(paths[0], LCF_ALLOWED_TYPES)
    if errors:
        raise ValueError("\n".join(errors))

    p("⚙️  Computing weekly demand…")
    weekly_df, util_df, types, years = core.build_weekly(df, col_map, capacities)

    p("📗 Building workbook…")
    wb = openpyxl.Workbook()
    wb.remove(wb.active)

    p("📄 Sheet 1 — Weekly Planner…")
    core.write_weekly_planner(wb, weekly_df, types, TOOL_TITLE)

    p("🎨 Sheet 2 — Utilization…")
    core.write_utilization_sheet(wb, util_df, types, theme)

    p("📋 Sheet 3 — Summary…")
    cap_list = " | ".join(f"{k}={v}/yr" for k,v in capacities.items())
    core.write_summary_sheet(wb, weekly_df, util_df, types, capacities, years,
                              TOOL_TITLE, cap_note=cap_list)

    p("📈 Sheet 4 — Utilization Chart…")
    core.write_utilization_chart(wb, util_df, types, years)

    p("📊 Sheet 5 — Capacity vs Demand Chart…")
    core.write_capacity_chart(wb, weekly_df, types, capacities, years)

    p("📊 Sheet 6 — Year-on-Year…")
    core.write_yoy_chart(wb, weekly_df, types, years)

    p("🗓️  Sheet 7 — Gantt All Years…")
    core.write_gantt_all_years(wb, weekly_df, types, capacities, years)

    p("🔥 Sheet 8 — Gantt Heatmap…")
    core.write_gantt_heatmap(wb, weekly_df, types, capacities, years)

    p("💾 Saving…")
    out = core.save_workbook(wb, paths[0], LCF_OUTPUT_FILE)
    return out, warns


# ═══════════════════════════════════════════════════════════
#  GUI
# ═══════════════════════════════════════════════════════════


# ============================================================
# generate_coating
# ============================================================

"""
Tool 2 — Coating Labs Dashboard (Cold Spray, HVOF, Plasma)
============================================================
Combined capacity = 350 samples/year across all 3 coating labs.
Individual + combined utilization tracking.
Run: python tool2_coating_labs.py
"""

import os, sys, threading, math


import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.chart.series import SeriesLabel
from openpyxl.utils import get_column_letter




# ── Tool-specific config ──────────────────────────────────
TOOL_TITLE    = "Coating Labs Dashboard (Cold Spray, HVOF, Plasma)"
COATING_ALLOWED_TYPES = ["Cold Spray", "HVOF", "Plasma"]   # unique name
COATING_OUTPUT_FILE = "Coating_Labs_Dashboard.xlsx"   # unique name — see note
                                                       # above LCF_OUTPUT_FILE
COMBINED_CAP  = 350   # total samples/year across all 3 coating labs

# Individual capacities for single-lab utilization display
# These sum to 350 (proportional split based on typical usage)
DEFAULT_CAPS = {
    "Cold Spray": 140,  # largest portion
    "HVOF":       120,
    "Plasma":      90,
}


# ═══════════════════════════════════════════════════════════
#  COMBINED UTILIZATION SHEET
# ═══════════════════════════════════════════════════════════

def write_combined_sheet(wb, weekly_df, types, years, combined_cap):
    """Shows total combined demand across all coating labs vs 350 cap."""
    ws = wb.create_sheet("Combined_Utilization")
    core.banner(ws, 1,
                f"Combined Coating Labs — Total Demand vs {combined_cap}/yr Capacity",
                cols=20)

    weekly_cap = combined_cap / 52

    # Build combined weekly
    headers = ["Year", "Week", "Total Demand", "Capacity/Wk", "Utilization %",
               "Status"] + types
    for c, h in enumerate(headers, 1):
        core.hdr(ws.cell(3, c, h), bg="1F3864", sz=9)
        core.cw(ws, c, 14)
    core.cw(ws, 1, 8); core.cw(ws, 2, 7); core.cw(ws, 3, 14)
    core.cw(ws, 4, 12); core.cw(ws, 5, 14); core.cw(ws, 6, 12)

    red_f  = PatternFill("solid", start_color="FF4444")
    yel_f  = PatternFill("solid", start_color="FFD700")
    grn_f  = PatternFill("solid", start_color="70AD47")

    for r, row in enumerate(weekly_df.itertuples(index=False), 4):
        vals     = list(row)
        year_v   = vals[0]; week_v = vals[1]
        type_vals = {t: vals[2 + i] for i, t in enumerate(types)}
        total    = sum(type_vals.values())
        util     = total / weekly_cap if weekly_cap > 0 else 0.0

        ws.cell(r, 1, year_v); core.fmt(ws.cell(r, 1))
        ws.cell(r, 2, week_v); core.fmt(ws.cell(r, 2))

        tc = ws.cell(r, 3, round(total, 3)); core.fmt(tc, number_format="0.000")
        if util > 1.0:   tc.fill = red_f;  tc.font = Font(name="Arial", size=9, bold=True, color="FFFFFF")
        elif util >= 0.8: tc.fill = yel_f
        else:             tc.fill = grn_f;  tc.font = Font(name="Arial", size=9, color="FFFFFF")

        ws.cell(r, 4, round(weekly_cap, 2)); core.fmt(ws.cell(r, 4), number_format="0.00")

        uc = ws.cell(r, 5, round(util, 4)); core.fmt(uc, number_format="0.0%")
        if util > 1.0:   uc.fill = red_f;  uc.font = Font(name="Arial", size=9, bold=True, color="FFFFFF")
        elif util >= 0.8: uc.fill = yel_f
        else:             uc.fill = grn_f;  uc.font = Font(name="Arial", size=9, color="FFFFFF")

        status = "🔴 OVERLOAD" if util > 1.0 else "🟡 Near Cap" if util >= 0.8 else "🟢 OK"
        ws.cell(r, 6, status); core.fmt(ws.cell(r, 6))

        for ci, t in enumerate(types, 7):
            ws.cell(r, ci, round(type_vals[t], 3)); core.fmt(ws.cell(r, ci), number_format="0.000")

    ws.freeze_panes = "C4"
    ws.auto_filter.ref = f"A3:{get_column_letter(len(headers))}3"


def write_combined_chart(wb, weekly_df, types, years, combined_cap):
    """Line chart: combined weekly demand vs combined capacity."""
    ws = wb.create_sheet("Chart_Combined")
    core.banner(ws, 1,
                f"Combined Coating Labs — Weekly Demand Trend vs {combined_cap}/yr Cap",
                cols=30)

    # Build monthly aggregated combined demand
    months = ["Jan","Feb","Mar","Apr","May","Jun",
              "Jul","Aug","Sep","Oct","Nov","Dec"]
    weekly_cap = combined_cap / 52

    HR = 3
    ws.cell(HR, 1, "Month").font = Font(bold=True, name="Arial", size=9)
    ws.cell(HR, 1).fill = PatternFill("solid", start_color="D9E1F2")

    yr_cols = {}
    for i, y in enumerate(years):
        col = 2 + i
        h = ws.cell(HR, col, str(y))
        h.font = Font(bold=True, name="Arial", size=9)
        h.fill = PatternFill("solid", start_color="D9E1F2")
        h.alignment = Alignment(horizontal="center")
        core.cw(ws, col, 14)
        yr_cols[y] = col
    core.cw(ws, 1, 10)

    # Capacity line column
    cap_col = 2 + len(years)
    hc = ws.cell(HR, cap_col, f"Cap/Wk ({combined_cap}/yr)")
    hc.font = Font(bold=True, name="Arial", size=9)
    hc.fill = PatternFill("solid", start_color="FFD700")
    hc.alignment = Alignment(horizontal="center")
    core.cw(ws, cap_col, 16)

    for mi, month in enumerate(months):
        row = HR + 1 + mi
        ws.cell(row, 1, month).font = Font(bold=True, name="Arial", size=9)
        ws.cell(row, 1).fill = PatternFill("solid", start_color="EBF3FB")

        for y in years:
            yd = weekly_df[weekly_df["Year"] == y].copy()
            yd["_m"] = yd["Week"].apply(core.week_to_month)
            month_rows = yd[yd["_m"] == mi + 1]
            total_avg  = sum(month_rows[t].mean() for t in types
                             if not month_rows.empty)
            if math.isnan(total_avg): total_avg = 0.0
            cell = ws.cell(row, yr_cols[y], round(total_avg, 3))
            cell.number_format = "0.000"
            cell.font      = Font(name="Arial", size=9)
            cell.alignment = Alignment(horizontal="center")

        # Capacity reference
        cc = ws.cell(row, cap_col, round(weekly_cap, 2))
        cc.number_format = "0.00"
        cc.font      = Font(name="Arial", size=9, bold=True)
        cc.fill      = PatternFill("solid", start_color="FFF2CC")
        cc.alignment = Alignment(horizontal="center")

    # Line chart
    chart = LineChart()
    chart.title        = f"Combined Coating Labs Weekly Demand (Cap={combined_cap}/yr)"
    chart.y_axis.title = "Weekly Demand (samples)"
    chart.x_axis.title = "Month"
    chart.style        = 10
    chart.height       = 16
    chart.width        = 30

    # Add each year series
    first_col = min(yr_cols.values())
    first_ref = Reference(ws, min_col=first_col, max_col=first_col,
                          min_row=HR, max_row=HR + 12)
    chart.add_data(first_ref, titles_from_data=True)
    cats = Reference(ws, min_col=1, min_row=HR+1, max_row=HR+12)
    chart.set_categories(cats)

    for y in list(years)[1:]:
        col = yr_cols[y]
        ref = Reference(ws, min_col=col, max_col=col, min_row=HR, max_row=HR+12)
        chart.add_data(ref, titles_from_data=True)

    # Capacity reference line
    cap_ref = Reference(ws, min_col=cap_col, max_col=cap_col,
                        min_row=HR, max_row=HR+12)
    chart.add_data(cap_ref, titles_from_data=True)

    ws.add_chart(chart, f"A{HR + 15}")

    # Also add stacked bar chart per type
    chart2 = BarChart()
    chart2.type     = "col"
    chart2.grouping = "stacked"
    chart2.title    = "Demand by Type — Stacked (Monthly)"
    chart2.y_axis.title = "Weekly Demand"
    chart2.x_axis.title = "Month"
    chart2.style    = 10
    chart2.height   = 16
    chart2.width    = 28

    # Write per-type monthly data for last year
    last_yr = max(years)
    type_start_col = cap_col + 2
    ws.cell(HR, type_start_col - 1, f"Per-Type {last_yr}").font = \
        Font(bold=True, name="Arial", size=9)
    for ti, t in enumerate(types):
        col = type_start_col + ti
        h2  = ws.cell(HR, col, t)
        h2.font = Font(bold=True, name="Arial", size=9)
        h2.fill = PatternFill("solid", start_color="D9E1F2")
        h2.alignment = Alignment(horizontal="center")
        core.cw(ws, col, 13)
        yd = weekly_df[weekly_df["Year"] == last_yr].copy()
        yd["_m"] = yd["Week"].apply(core.week_to_month)
        for mi in range(12):
            row  = HR + 1 + mi
            grp  = yd[yd["_m"] == mi + 1][t]
            avg  = grp.mean() if not grp.empty else 0.0
            cell = ws.cell(row, col, round(avg if not math.isnan(avg) else 0.0, 3))
            cell.number_format = "0.000"
            cell.alignment     = Alignment(horizontal="center")
            cell.font          = Font(name="Arial", size=9)

    for ti, t in enumerate(types):
        col  = type_start_col + ti
        ref2 = Reference(ws, min_col=col, max_col=col, min_row=HR, max_row=HR+12)
        chart2.add_data(ref2, titles_from_data=True)
    chart2.set_categories(cats)
    chart2_col = get_column_letter(type_start_col)
    ws.add_chart(chart2, f"{chart2_col}{HR + 15}")


# ═══════════════════════════════════════════════════════════
#  MAIN GENERATOR
# ═══════════════════════════════════════════════════════════

def generate_coating(input_path, individual_caps, theme, progress_cb=None):
    def p(msg):
        if progress_cb: progress_cb(msg)

    p("📖 Reading file…")
    paths = input_path if isinstance(input_path, (list, tuple)) else [input_path]
    if len(paths) > 1:
        df, col_map, errors, warns = merge_files(paths, allowed_types=COATING_ALLOWED_TYPES)
    else:
        df, col_map, errors, warns = core.load_and_filter(paths[0], COATING_ALLOWED_TYPES)
    if errors:
        raise ValueError("\n".join(errors))

    p("⚙️  Computing weekly demand…")
    weekly_df, util_df, types, years = core.build_weekly(df, col_map, individual_caps)

    p("📗 Building workbook…")
    wb = openpyxl.Workbook()
    wb.remove(wb.active)

    p("📄 Sheet 1 — Weekly Planner…")
    core.write_weekly_planner(wb, weekly_df, types, TOOL_TITLE)

    p("🎨 Sheet 2 — Utilization…")
    core.write_utilization_sheet(wb, util_df, types, theme)

    p("📋 Sheet 3 — Summary…")
    cap_note2 = " | ".join(f"{k}={v}/yr" for k,v in individual_caps.items())
    core.write_summary_sheet(wb, weekly_df, util_df, types, individual_caps, years,
                              TOOL_TITLE, cap_note=cap_note2)

    p("📈 Sheet 4 — Utilization Chart…")
    core.write_utilization_chart(wb, util_df, types, years)

    p("📊 Sheet 5 — Capacity vs Demand Chart…")
    core.write_capacity_chart(wb, weekly_df, types, individual_caps, years)

    p("📊 Sheet 6 — Year-on-Year…")
    core.write_yoy_chart(wb, weekly_df, types, years)

    p("🗓️  Sheet 7 — Gantt All Years…")
    core.write_gantt_all_years(wb, weekly_df, types, individual_caps, years)

    p("🔥 Sheet 8 — Gantt Heatmap…")
    core.write_gantt_heatmap(wb, weekly_df, types, individual_caps, years)

    p("💾 Saving…")
    out = core.save_workbook(wb, paths[0], COATING_OUTPUT_FILE)
    return out, warns


# ────────────────────────────────────────────────────────────
#  GUI
# ────────────────────────────────────────────────────────────


# ============================================================
# generate_thermal
# ============================================================

"""
Tool 3 — Thermal Lab Dashboard
================================
Tracks up to N Thermal Rigs (dynamic — rig names driven by rig_caps dict).
Default rigs: Thermal Rig 1, Thermal Rig 2, Thermal Rig 3 (20 samples/yr each).
"""

THERMAL_TOOL_TITLE = "Thermal Lab Dashboard"
THERMAL_OUTPUT_FILE = "Thermal_Lab_Dashboard.xlsx"


def generate_thermal(input_path, rig_caps, theme, progress_cb=None):
    """
    Generate an 8-sheet Thermal Lab Excel dashboard.

    Parameters
    ----------
    input_path  : str   — path to the user's Excel file
    rig_caps    : dict  — {rig_name: annual_capacity}, e.g.
                          {"Thermal Rig 1": 20, "Thermal Rig 2": 20, ...}
    theme       : dict  — core.COLOR_THEMES entry
    progress_cb : callable or None
    """
    def p(msg):
        if progress_cb:
            progress_cb(msg)

    allowed_types = list(rig_caps.keys())

    p("📖 Reading file…")
    paths = input_path if isinstance(input_path, (list, tuple)) else [input_path]
    if len(paths) > 1:
        df, col_map, errors, warns = merge_files(paths, allowed_types=allowed_types)
    else:
        df, col_map, errors, warns = core.load_and_filter(paths[0], allowed_types)
    if errors:
        raise ValueError("\n".join(errors))

    p("⚙️  Computing weekly demand…")
    weekly_df, util_df, types, years = core.build_weekly(df, col_map, rig_caps)

    p("📗 Building workbook…")
    wb = openpyxl.Workbook()
    wb.remove(wb.active)

    p("📄 Sheet 1 — Weekly Planner…")
    core.write_weekly_planner(wb, weekly_df, types, THERMAL_TOOL_TITLE)

    p("🎨 Sheet 2 — Utilization…")
    core.write_utilization_sheet(wb, util_df, types, theme)

    p("📋 Sheet 3 — Summary…")
    cap_list = " | ".join(f"{k}={v}/yr" for k, v in rig_caps.items())
    core.write_summary_sheet(wb, weekly_df, util_df, types, rig_caps, years,
                              THERMAL_TOOL_TITLE, cap_note=cap_list)

    p("📈 Sheet 4 — Utilization Chart…")
    core.write_utilization_chart(wb, util_df, types, years)

    p("📊 Sheet 5 — Capacity vs Demand Chart…")
    core.write_capacity_chart(wb, weekly_df, types, rig_caps, years)

    p("📊 Sheet 6 — Year-on-Year…")
    core.write_yoy_chart(wb, weekly_df, types, years)

    p("🗓️  Sheet 7 — Gantt All Years…")
    core.write_gantt_all_years(wb, weekly_df, types, rig_caps, years)

    p("🔥 Sheet 8 — Gantt Heatmap…")
    core.write_gantt_heatmap(wb, weekly_df, types, rig_caps, years)

    p("💾 Saving…")
    out = core.save_workbook(wb, paths[0], THERMAL_OUTPUT_FILE)
    return out, warns


# ============================================================
# generate_comparison
# ============================================================

"""
Tool 3 — Lab Comparison Dashboard  (Multi-File Edition)
=========================================================
Supports ONE or MULTIPLE Excel files — data is merged automatically.
  Mode A: Single file with all lab types
  Mode B: Two files (one per group)
  Mode C: Multiple files — all merged by lab type + year

Run: python tool3_comparison.py
"""

import os, sys, threading, math


import pandas as pd
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.chart import BarChart, LineChart, Reference
from openpyxl.chart.series import SeriesLabel
from openpyxl.utils import get_column_letter




TOOL_TITLE  = "Lab Comparison Dashboard — Mechanical vs Coating vs Thermal Labs"
OUTPUT_FILE = "Lab_Comparison_Dashboard.xlsx"

GROUP_A = {"name":"Mechanical Labs","types":["LCF","Creep"],
           "caps":{"LCF":50,"Creep":22},"color":"1A4E8A","cap_total":72}
GROUP_B = {"name":"Coating Labs","types":["Cold Spray","HVOF","Plasma"],
           "caps":{"Cold Spray":140,"HVOF":120,"Plasma":90},"color":"1F5C1A","cap_total":350}
GROUP_C = {"name":"Thermal Lab","types":["Thermal Rig"],
           "caps":{"Thermal Rig":20},
           "color":"7F3F00","cap_total":20}
ALL_TYPES = GROUP_A["types"] + GROUP_B["types"] + GROUP_C["types"]


# ────────────────────────────────────────────────────────────
#  MULTI-FILE MERGE ENGINE
# ────────────────────────────────────────────────────────────

def merge_files(file_paths, allowed_types=None):
    """
    Load >=1 Excel files using the smart parser (core.load_and_filter), normalise,
    and merge into a single long-format DataFrame.
    Duplicate Year+Type rows across files are SUMMED.

    allowed_types restricts which lab types are accepted from each file —
    pass a tool's own type list (e.g. ["LCF","Creep"]) so that tool only
    ever looks at its own assigned task, even if a file happens to contain
    other lab types' data too. Defaults to ALL_TYPES (the original Tool 4
    behaviour, which needs to recognise any of the 6 lab types from any
    file — including messy/monthly-block formats handled by load_and_filter).

    Returns: (merged_df, col_map, errors, file_log)
    """
    if allowed_types is None:
        allowed_types = ALL_TYPES

    frames = []
    errors = []
    file_log = []   # one line per file for the Data_Sources sheet

    for path in file_paths:
        fname = os.path.basename(path)
        try:
            df_loaded, col_map, errs_l, warns_l = core.load_and_filter(path, allowed_types)
        except Exception as e:
            errors.append(f"Cannot process '{fname}': {e}")
            continue

        if errs_l:
            errors.append(f"'{fname}': {'; '.join(errs_l)}")
            continue

        if df_loaded is None or df_loaded.empty:
            file_log.append(f"SKIPPED (empty/no matching types): {fname}")
            continue

        yc = col_map.get("year", "Year")
        tc = col_map.get("type", "Type")
        vc = col_map.get("value", "Value")

        df = df_loaded[[yc, tc, vc]].copy()
        df.columns = ["Year", "Type", "Value"]
        df["Value"] = pd.to_numeric(df["Value"], errors="coerce").fillna(0)
        df["Year"]  = pd.to_numeric(df["Year"],  errors="coerce")
        df = df.dropna(subset=["Year","Type","Value"])
        df["Year"]  = df["Year"].astype(int)
        df["Type"]  = df["Type"].astype(str).str.strip()

        # Normalise type names (case-insensitive) to canonical spelling
        known = {t.lower(): t for t in allowed_types}
        df["Type"] = df["Type"].apply(lambda x: known.get(x.lower(), x))

        rows_kept = len(df)
        if rows_kept == 0:
            file_log.append(f"SKIPPED (no usable rows): {fname}")
            continue

        frames.append(df)
        log_line = f"OK ({rows_kept} rows): {fname}"
        if warns_l:
            log_line += f"  [{'; '.join(warns_l)}]"
        file_log.append(log_line)

    if errors:
        return None, {}, errors, file_log

    if not frames:
        errors.append("No valid data found in any file.")
        return None, {}, errors, file_log

    merged = pd.concat(frames, ignore_index=True)
    # Sum duplicates (same Year+Type from multiple files)
    merged = merged.groupby(["Year","Type"], as_index=False)["Value"].sum()

    found   = set(merged["Type"].unique())
    missing = set(allowed_types) - found
    if missing:
        file_log.append(f"WARNING: Lab types not found in any file: {sorted(missing)}")

    col_map_out = {"year":"Year","type":"Type","value":"Value"}
    return merged, col_map_out, [], file_log


# ────────────────────────────────────────────────────────────
#  SHEET WRITERS
# ────────────────────────────────────────────────────────────

def write_data_sources(wb, file_paths, file_log):
    ws = wb.create_sheet("Data_Sources")
    core.banner(ws, 1, "Data Sources Loaded", cols=12)
    for c,h in enumerate(["#","File Name","Status","Full Path"],1):
        core.hdr(ws.cell(3,c,h), bg="1F3864", sz=9)
    core.cw(ws,1,5); core.cw(ws,2,35); core.cw(ws,3,18); core.cw(ws,4,60)
    for i,path in enumerate(file_paths,4):
        log = file_log[i-4] if i-4 < len(file_log) else ""
        status = "✓ OK" if log.startswith("OK") else "⚠ WARNING" if "WARNING" in log else "✗ ERROR"
        ws.cell(i,1,i-3).border=core.bdr()
        ws.cell(i,1).font=Font(name="Arial",size=9)
        c2=ws.cell(i,2,os.path.basename(path))
        c2.font=Font(name="Arial",size=9,bold=True); c2.border=core.bdr()
        c2.alignment=Alignment(horizontal="left")
        c3=ws.cell(i,3,status)
        c3.font=Font(name="Arial",size=9); c3.border=core.bdr()
        c3.fill=PatternFill("solid",start_color="C6EFCE" if "OK" in status else "FFD700")
        c3.alignment=Alignment(horizontal="center")
        c4=ws.cell(i,4,path)
        c4.font=Font(name="Arial",size=8,color="666666"); c4.border=core.bdr()
        c4.alignment=Alignment(horizontal="left")


def write_overview(wb, wdf_a, wdf_b, years, wdf_c=None, caps_c=None):
    ws = wb.create_sheet("Overview")
    core.banner(ws, 1, f"Lab Comparison — Overview", cols=24, sz=13)
    R = 3
    groups = [(GROUP_A, wdf_a), (GROUP_B, wdf_b)]
    if wdf_c is not None:
        caps_c_total = sum(caps_c.values()) if caps_c else GROUP_C["cap_total"]
        grp_c = dict(GROUP_C)
        grp_c["cap_total"] = caps_c_total
        grp_c["caps"]      = caps_c or GROUP_C["caps"]
        grp_c["types"]     = list(caps_c.keys()) if caps_c else GROUP_C["types"]
        groups.append((grp_c, wdf_c))
    for grp, wdf in groups:
        ws.merge_cells(start_row=R,start_column=1,end_row=R,end_column=3+len(years))
        c=ws.cell(R,1,f"  {grp['name']}  —  Capacity: {grp['cap_total']} samples/yr")
        c.font=Font(bold=True,size=11,name="Arial",color="FFFFFF")
        c.fill=PatternFill("solid",start_color=grp["color"])
        c.alignment=Alignment(horizontal="left",vertical="center")
        ws.row_dimensions[R].height=20; R+=1
        tgt_bg="2E75B6" if grp==GROUP_A else "375623"
        for ci,h in enumerate(["Lab Type"]+[str(y) for y in years]+["Total"],1):
            core.hdr(ws.cell(R,ci,h),bg=tgt_bg,sz=9); core.cw(ws,ci,14)
        R+=1
        for t in grp["types"]:
            ws.cell(R,1,t).font=Font(name="Arial",size=9,bold=True)
            ws.cell(R,1).border=core.bdr()
            ws.cell(R,1).alignment=Alignment(horizontal="left")
            row_total=0
            for ci,y in enumerate(years,2):
                val=wdf[wdf["Year"]==y][t].sum() if t in wdf.columns else 0
                cell=ws.cell(R,ci,round(val,1)); core.fmt(cell,number_format="#,##0.0")
                row_total+=val
            tc=ws.cell(R,2+len(years),round(row_total,1))
            core.fmt(tc,number_format="#,##0.0",bold=True)
            tc.fill=PatternFill("solid",start_color="BDD7EE" if grp==GROUP_A else "C6EFCE")
            R+=1
        ws.cell(R,1,"GROUP TOTAL")
        ws.cell(R,1).fill=PatternFill("solid",start_color=grp["color"])
        ws.cell(R,1).font=Font(name="Arial",size=9,bold=True,color="FFFFFF")
        ws.cell(R,1).border=core.bdr()
        ws.cell(R,1).alignment=Alignment(horizontal="left")
        grand=0
        for ci,y in enumerate(years,2):
            yr_tot=sum(wdf[wdf["Year"]==y][t].sum() for t in grp["types"] if t in wdf.columns)
            cell=ws.cell(R,ci,round(yr_tot,1)); core.fmt(cell,number_format="#,##0.0",bold=True)
            util=yr_tot/grp["cap_total"] if grp["cap_total"]>0 else 0
            cell.fill=PatternFill("solid",start_color="FF4444" if util>1 else "FFD700" if util>=0.8 else "C6EFCE")
            if util>1: cell.font=Font(name="Arial",size=9,bold=True,color="FFFFFF")
            grand+=yr_tot
        tc=ws.cell(R,2+len(years),round(grand,1)); core.fmt(tc,number_format="#,##0.0",bold=True)
        tc.fill=PatternFill("solid",start_color=grp["color"])
        tc.font=Font(name="Arial",size=9,bold=True,color="FFFFFF"); R+=3
    # Util summary
    core.section_hdr(ws,R,"Annual Utilization vs Capacity",3+len(years)); R+=1
    for ci,h in enumerate(["Lab Group","Cap/yr"]+[str(y) for y in years],1):
        core.hdr(ws.cell(R,ci,h),bg="C00000",sz=9); core.cw(ws,ci,16)
    R+=1
    for grp,wdf in groups:
        ws.cell(R,1,grp["name"])
        ws.cell(R,1).fill=PatternFill("solid",start_color=grp["color"])
        ws.cell(R,1).font=Font(name="Arial",size=9,bold=True,color="FFFFFF")
        ws.cell(R,1).border=core.bdr(); ws.cell(R,1).alignment=Alignment(horizontal="left")
        ws.cell(R,2,grp["cap_total"]); core.fmt(ws.cell(R,2),number_format="#,##0")
        for ci,y in enumerate(years,3):
            tot=sum(wdf[wdf["Year"]==y][t].sum() for t in grp["types"] if t in wdf.columns)
            util=tot/grp["cap_total"] if grp["cap_total"]>0 else 0
            cell=ws.cell(R,ci,round(util,4)); core.fmt(cell,number_format="0.0%")
            if util>1.0: cell.fill=PatternFill("solid",start_color="FF4444"); cell.font=Font(name="Arial",size=9,bold=True,color="FFFFFF")
            elif util>=0.8: cell.fill=PatternFill("solid",start_color="FFD700"); cell.font=Font(name="Arial",size=9,bold=True)
            else: cell.fill=PatternFill("solid",start_color="70AD47"); cell.font=Font(name="Arial",size=9,color="FFFFFF")
        R+=1


def write_comparison_chart(wb, wdf_a, wdf_b, years, wdf_c=None, caps_c=None):
    ws = wb.create_sheet("Chart_Comparison")
    core.banner(ws,1,"Group Comparison — Annual Demand vs Capacity",cols=24)
    months=["Jan","Feb","Mar","Apr","May","Jun","Jul","Aug","Sep","Oct","Nov","Dec"]
    HR=3; D0=4
    for ci,h in enumerate(["Year",f"{GROUP_A['name']} Demand",f"Cap({GROUP_A['cap_total']})",
                            f"{GROUP_B['name']} Demand",f"Cap({GROUP_B['cap_total']})"],1):
        core.hdr(ws.cell(HR,ci,h),bg="1F3864",sz=9); core.cw(ws,ci,22)
    DE=D0+len(years)-1
    for ri,y in enumerate(years,D0):
        tot_a=sum(wdf_a[wdf_a["Year"]==y][t].sum() for t in GROUP_A["types"] if t in wdf_a.columns)
        tot_b=sum(wdf_b[wdf_b["Year"]==y][t].sum() for t in GROUP_B["types"] if t in wdf_b.columns)
        ws.cell(ri,1,y); core.fmt(ws.cell(ri,1))
        for ci,tot,cap in [(2,tot_a,GROUP_A["cap_total"]),(4,tot_b,GROUP_B["cap_total"])]:
            cell=ws.cell(ri,ci,round(tot,1)); core.fmt(cell,number_format="#,##0.0")
            util=tot/cap if cap>0 else 0
            cell.fill=PatternFill("solid",start_color="FF4444" if util>1 else "FFD700" if util>=0.8 else "C6EFCE")
            if util>1: cell.font=Font(name="Arial",size=9,bold=True,color="FFFFFF")
        ws.cell(ri,3,GROUP_A["cap_total"]); core.fmt(ws.cell(ri,3),number_format="#,##0")
        ws.cell(ri,5,GROUP_B["cap_total"]); core.fmt(ws.cell(ri,5),number_format="#,##0")
    chart=BarChart(); chart.type="col"; chart.grouping="clustered"
    chart.title="Annual Demand — Mechanical vs Coating Labs"
    chart.y_axis.title="Samples/Year"; chart.x_axis.title="Year"
    chart.style=10; chart.height=15; chart.width=28
    cats=Reference(ws,min_col=1,min_row=D0,max_row=DE)
    first=Reference(ws,min_col=2,max_col=2,min_row=HR,max_row=DE)
    chart.add_data(first,titles_from_data=True); chart.set_categories(cats)
    for ci in [4,3,5]:
        ref=Reference(ws,min_col=ci,max_col=ci,min_row=HR,max_row=DE)
        chart.add_data(ref,titles_from_data=True)
    ws.add_chart(chart,f"A{DE+4}")
    # Monthly util line
    UR=DE+22
    core.section_hdr(ws,UR,"Monthly Utilization % (last 2 years)",10); UR+=1
    ws.cell(UR,1,"Month").font=Font(bold=True,name="Arial",size=9)
    ws.cell(UR,1).fill=PatternFill("solid",start_color="D9E1F2")
    plot_years=list(years[-2:]) if len(years)>=2 else list(years)
    col_a={y:2+2*i for i,y in enumerate(plot_years)}
    col_b={y:3+2*i for i,y in enumerate(plot_years)}
    for y in plot_years:
        for col,label in [(col_a[y],f"Mech {y}"),(col_b[y],f"Coating {y}")]:
            h=ws.cell(UR,col,label); h.font=Font(bold=True,name="Arial",size=8)
            h.fill=PatternFill("solid",start_color="D9E1F2")
            h.alignment=Alignment(horizontal="center"); core.cw(ws,col,14)
    core.cw(ws,1,10)
    for mi,month in enumerate(months):
        row=UR+1+mi
        ws.cell(row,1,month).font=Font(bold=True,name="Arial",size=9)
        ws.cell(row,1).fill=PatternFill("solid",start_color="EBF3FB")
        for y in plot_years:
            for wdf,grp,col in [(wdf_a,GROUP_A,col_a[y]),(wdf_b,GROUP_B,col_b[y])]:
                yd=wdf[wdf["Year"]==y].copy(); yd["_m"]=yd["Week"].apply(core.week_to_month)
                tot=sum(yd[yd["_m"]==mi+1][t].mean() for t in grp["types"] if t in yd.columns)
                util=tot/(grp["cap_total"]/52) if grp["cap_total"]>0 else 0
                if math.isnan(util): util=0.0
                cell=ws.cell(row,col,round(util,4)); cell.number_format="0.0%"
                cell.font=Font(name="Arial",size=9); cell.alignment=Alignment(horizontal="center")
    chart2=LineChart()
    chart2.title="Monthly Utilization — Mechanical vs Coating Labs"
    chart2.y_axis.title="Utilization %"; chart2.x_axis.title="Month"
    chart2.y_axis.numFmt="0%"; chart2.y_axis.scaling.min=0
    chart2.style=10; chart2.height=14; chart2.width=28
    all_cols=list(col_a.values())+list(col_b.values())
    first_col=min(all_cols)
    first_ref=Reference(ws,min_col=first_col,max_col=first_col,min_row=UR,max_row=UR+12)
    chart2.add_data(first_ref,titles_from_data=True)
    cats2=Reference(ws,min_col=1,min_row=UR+1,max_row=UR+12); chart2.set_categories(cats2)
    for col in all_cols:
        if col==first_col: continue
        ref=Reference(ws,min_col=col,max_col=col,min_row=UR,max_row=UR+12)
        chart2.add_data(ref,titles_from_data=True)
    ws.add_chart(chart2,f"A{UR+15}")


def write_all_labs(wb, wdf_a, wdf_b, years, types_a, types_b,
                   wdf_c=None, types_c=None, caps_c=None):
    ws=wb.create_sheet("All_Labs_Summary")
    core.banner(ws,1,"All Labs — Combined Annual Demand Summary",cols=20)
    all_types=types_a+types_b+(list(types_c) if types_c else [])
    all_wdfs={t:wdf_a for t in types_a}; all_wdfs.update({t:wdf_b for t in types_b})
    if wdf_c is not None and types_c:
        all_wdfs.update({t:wdf_c for t in types_c})
    all_caps={**GROUP_A["caps"],**GROUP_B["caps"]}
    if caps_c: all_caps.update(caps_c)
    else: all_caps.update(GROUP_C["caps"])
    HR=3; D0=4
    types_c_local = list(types_c) if types_c else []
    n_a, n_b = len(types_a), len(types_b)
    hdrs=["Year"]+all_types+["Mech Total","Coating Total","Thermal Total","Grand Total"]
    for ci,h in enumerate(hdrs,1):
        if ci==1: bg="1F3864"
        elif ci<=1+n_a: bg=GROUP_A["color"]
        elif ci<=1+n_a+n_b: bg=GROUP_B["color"]
        elif ci<=1+len(all_types): bg=GROUP_C["color"]
        else: bg="C00000"
        core.hdr(ws.cell(HR,ci,h),bg=bg,sz=9); core.cw(ws,ci,13)
    for ri,y in enumerate(years,D0):
        ws.cell(ri,1,y); core.fmt(ws.cell(ri,1))
        mech_tot=coat_tot=therm_tot=0
        for ci,t in enumerate(all_types,2):
            wdf=all_wdfs[t]; val=wdf[wdf["Year"]==y][t].sum() if t in wdf.columns else 0
            cell=ws.cell(ri,ci,round(val,1)); core.fmt(cell,number_format="#,##0.0")
            cap=all_caps.get(t,1); util=val/cap if cap>0 else 0
            cell.fill=PatternFill("solid",start_color="FF4444" if util>1 else "FFD700" if util>=0.8 else "C6EFCE")
            if util>1: cell.font=Font(name="Arial",size=9,bold=True,color="FFFFFF")
            if t in types_a: mech_tot+=val
            elif t in types_b: coat_tot+=val
            else: therm_tot+=val
        lc=1+len(all_types)
        mc=ws.cell(ri,lc+1,round(mech_tot,1)); core.fmt(mc,number_format="#,##0.0",bold=True); mc.fill=PatternFill("solid",start_color="BDD7EE")
        cc=ws.cell(ri,lc+2,round(coat_tot,1)); core.fmt(cc,number_format="#,##0.0",bold=True); cc.fill=PatternFill("solid",start_color="C6EFCE")
        tc2=ws.cell(ri,lc+3,round(therm_tot,1)); core.fmt(tc2,number_format="#,##0.0",bold=True); tc2.fill=PatternFill("solid",start_color="FCE4D6")
        gc=ws.cell(ri,lc+4,round(mech_tot+coat_tot+therm_tot,1)); core.fmt(gc,number_format="#,##0.0",bold=True); gc.fill=PatternFill("solid",start_color="D9D9D9")
    chart=BarChart(); chart.type="col"; chart.grouping="clustered"
    chart.title="All Labs — Annual Demand"; chart.y_axis.title="Samples/Year"
    chart.style=10; chart.height=16; chart.width=34
    cats=Reference(ws,min_col=1,min_row=D0,max_row=D0+len(years)-1)
    first=Reference(ws,min_col=2,max_col=2,min_row=HR,max_row=D0+len(years)-1)
    chart.add_data(first,titles_from_data=True); chart.set_categories(cats)
    for ci in range(3,2+len(all_types)):
        ref=Reference(ws,min_col=ci,max_col=ci,min_row=HR,max_row=D0+len(years)-1)
        chart.add_data(ref,titles_from_data=True)
    ws.add_chart(chart,f"A{D0+len(years)+4}")


# ────────────────────────────────────────────────────────────
#  MAIN GENERATOR
# ────────────────────────────────────────────────────────────

def generate_comparison(file_paths, caps_a, caps_b, theme, progress_cb=None,
                        caps_c=None):
    """
    Generate 3-group comparison Excel dashboard.
    caps_c: dict of {rig_name: capacity} for Thermal group (optional – defaults to GROUP_C caps).
    """
    def p(msg):
        if progress_cb: progress_cb(msg)

    if caps_c is None:
        caps_c = dict(GROUP_C["caps"])

    p(f"Loading {len(file_paths)} file(s)…")
    merged_df, col_map, errors, file_log = merge_files(file_paths)
    if errors:
        raise ValueError("\n".join(errors))

    tc = col_map["type"]
    df_a = merged_df[merged_df[tc].isin(GROUP_A["types"])].copy()
    df_b = merged_df[merged_df[tc].isin(GROUP_B["types"])].copy()
    df_c = merged_df[merged_df[tc].isin(list(caps_c.keys()))].copy()

    if df_a.empty:
        raise ValueError(f"No data for Mechanical labs {GROUP_A['types']}.")
    if df_b.empty:
        raise ValueError(f"No data for Coating labs {GROUP_B['types']}.")
    # Thermal group is optional — warn but continue
    has_c = not df_c.empty

    p("Computing weekly data…")
    wdf_a, udf_a, types_a, years_a = core.build_weekly(df_a, col_map, caps_a)
    wdf_b, udf_b, types_b, years_b = core.build_weekly(df_b, col_map, caps_b)
    years = sorted(set(list(years_a)) | set(list(years_b)))

    if has_c:
        wdf_c, udf_c, types_c, years_c = core.build_weekly(df_c, col_map, caps_c)
        years = sorted(set(years) | set(list(years_c)))
    else:
        wdf_c = udf_c = types_c = years_c = None
        file_log.append("WARNING: No Thermal lab data found — Group C omitted from Excel.")

    p("Building workbook…")
    wb = openpyxl.Workbook(); wb.remove(wb.active)

    p("Sheet 1 — Data Sources…")
    write_data_sources(wb, file_paths, file_log)
    p("Sheet 2 — Overview…")
    write_overview(wb, wdf_a, wdf_b, years, wdf_c=wdf_c, caps_c=caps_c)
    if has_c:
        p("Sheet 3 — All Labs Summary…")
        write_all_labs(wb, wdf_a, wdf_b, years, types_a, types_b,
                       wdf_c=wdf_c, types_c=types_c, caps_c=caps_c)
    p("Sheet 4 — Comparison Charts…")
    write_comparison_chart(wb, wdf_a, wdf_b, years,
                           wdf_c=wdf_c, caps_c=caps_c)
    p("Sheet 5 — Mech Utilization…")
    core.write_utilization_sheet(wb, udf_a, types_a, theme)
    wb.worksheets[-1].title = "Mech_Utilization"
    p("Sheet 6 — Coating Utilization…")
    core.write_utilization_sheet(wb, udf_b, types_b, theme)
    wb.worksheets[-1].title = "Coating_Utilization"
    if has_c:
        p("Sheet 7 — Thermal Utilization…")
        core.write_utilization_sheet(wb, udf_c, types_c, theme)
        wb.worksheets[-1].title = "Thermal_Utilization"
    p("Sheet 8 — Mech YoY…")
    core.write_yoy_chart(wb, wdf_a, types_a, years_a)
    wb.worksheets[-1].title = "Mech_YoY"
    p("Sheet 9 — Coating YoY…")
    core.write_yoy_chart(wb, wdf_b, types_b, years_b)
    wb.worksheets[-1].title = "Coating_YoY"
    if has_c:
        p("Sheet 10 — Thermal YoY…")
        core.write_yoy_chart(wb, wdf_c, types_c, years_c)
        wb.worksheets[-1].title = "Thermal_YoY"

    p("Saving…")
    out = core.save_workbook(wb, file_paths[0], OUTPUT_FILE)
    return out, file_log


# ═══════════════════════════════════════════════════════════
#  POWERPOINT EXPORT
# ═══════════════════════════════════════════════════════════

def generate_comparison_ppt(
    wdf_a, wdf_b, wdf_c,
    annual_a, annual_b, annual_c,
    types_a, types_b, types_c,
    caps_a, caps_b, caps_c,
    years,
    fig_cmp, fig_util, figs_gantt,
):
    """
    Build a 6-slide PowerPoint report and return it as bytes.

    Parameters
    ----------
    wdf_a/b/c    : weekly demand DataFrames per group
    annual_a/b/c : dict {year: {type: total}} per group
    types_a/b/c  : list of lab type names per group
    caps_a/b/c   : dict {type: annual_cap} per group
    years        : sorted list of all years
    fig_cmp      : Plotly Figure — comparison bar chart (Tab 1)
    fig_util     : Plotly Figure — utilisation line chart (Tab 3)
    figs_gantt   : list of Plotly Figures [gantt_a, gantt_b, gantt_c]

    Returns
    -------
    bytes  — the .pptx file contents
    """
    import io as _io
    import datetime
    import math
    import warnings
    warnings.filterwarnings("ignore")

    # ── optional imports (graceful fallback) ──────────────
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError:
        raise ImportError("python-pptx is required. Install: pip install python-pptx")

    try:
        import plotly.io as pio
        _HAS_KALEIDO = True
    except Exception:
        _HAS_KALEIDO = False

    # ── helpers ───────────────────────────────────────────
    W  = Inches(13.33)   # 16:9 widescreen
    H  = Inches(7.5)

    DARK_BLUE  = RGBColor(0x1F, 0x38, 0x64)
    MED_BLUE   = RGBColor(0x2E, 0x75, 0xB6)
    PURPLE     = RGBColor(0x7B, 0x5E, 0xA7)
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
    RED        = RGBColor(0xFF, 0x44, 0x44)
    AMBER      = RGBColor(0xFF, 0xD7, 0x00)
    GREEN      = RGBColor(0x70, 0xAD, 0x47)
    LIGHT_GREY = RGBColor(0xF0, 0xF4, 0xF8)

    GROUP_COLORS = {
        "Mechanical": RGBColor(0x1A, 0x4E, 0x8A),
        "Coating":    RGBColor(0x1F, 0x5C, 0x1A),
        "Thermal":    RGBColor(0x7F, 0x3F, 0x00),
    }

    def _solid_bg(slide, rgb):
        """Fill slide background with a solid colour."""
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb

    def _add_textbox(slide, text, left, top, width, height,
                     font_size=18, bold=False, color=WHITE,
                     align=PP_ALIGN.LEFT, wrap=True):
        txb = slide.shapes.add_textbox(left, top, width, height)
        tf  = txb.text_frame
        tf.word_wrap = wrap
        p   = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(font_size)
        run.font.bold  = bold
        run.font.color.rgb = color
        return txb

    def _add_title_bar(slide, title_text, subtitle_text=""):
        """Dark blue title bar at top of content slides."""
        bar = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            Inches(0), Inches(0), W, Inches(1.1)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = DARK_BLUE
        bar.line.fill.background()
        _add_textbox(slide, title_text,
                     Inches(0.25), Inches(0.05), Inches(12.5), Inches(0.65),
                     font_size=24, bold=True, color=WHITE)
        if subtitle_text:
            _add_textbox(slide, subtitle_text,
                         Inches(0.25), Inches(0.65), Inches(12.5), Inches(0.4),
                         font_size=13, color=RGBColor(0xBD, 0xD7, 0xEE))

    def _fig_to_png(fig, width=1200, height=600):
        """Convert Plotly figure to PNG bytes; returns None if kaleido unavailable."""
        if not _HAS_KALEIDO:
            return None
        try:
            import plotly.io as pio
            return pio.to_image(fig, format="png", width=width, height=height, scale=1.5)
        except Exception:
            return None

    def _embed_png(slide, png_bytes, left, top, width, height):
        """Embed PNG bytes as picture on slide."""
        buf = _io.BytesIO(png_bytes)
        slide.shapes.add_picture(buf, left, top, width, height)

    def _placeholder_box(slide, label, left, top, width, height):
        """Grey box shown when kaleido PNG is unavailable."""
        box = slide.shapes.add_shape(1, left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        box.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        _add_textbox(slide, f"[Chart: {label}]\n(kaleido not available)",
                     left + Inches(0.1), top + Inches(0.1),
                     width - Inches(0.2), height - Inches(0.2),
                     font_size=14, color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.CENTER)

    # ── build presentation ────────────────────────────────
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]   # completely blank layout

    last_year = int(max(years))
    today_str = datetime.date.today().strftime("%d %b %Y")

    # ── gather KPI data ───────────────────────────────────
    def _group_kpi(annual, caps, types, label):
        dem  = sum(annual.get(last_year, {}).get(t, 0) for t in types)
        cap  = sum(caps.values())
        util = dem / cap if cap > 0 else 0
        return {"label": label, "dem": dem, "cap": cap, "util": util}

    kpis = [
        _group_kpi(annual_a, caps_a, types_a, "🔵 Mechanical"),
        _group_kpi(annual_b, caps_b, types_b, "🟢 Coating"),
        _group_kpi(annual_c, caps_c, types_c, "🟠 Thermal"),
    ]

    # ══════════════════════════════════════════════════════
    #  SLIDE 1 — Title
    # ══════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(blank)
    _solid_bg(s1, DARK_BLUE)

    # purple accent bar left
    bar = s1.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.35), H)
    bar.fill.solid(); bar.fill.fore_color.rgb = PURPLE
    bar.line.fill.background()

    _add_textbox(s1, "Lab Planning &\nOccupancy Report",
                 Inches(0.6), Inches(1.6), Inches(11), Inches(2.5),
                 font_size=44, bold=True, color=WHITE)
    _add_textbox(s1,
                 f"Generated: {today_str}  |  Data up to {last_year}  |  "
                 f"All 3 Lab Groups: Mechanical · Coating · Thermal",
                 Inches(0.6), Inches(4.1), Inches(12), Inches(0.7),
                 font_size=16, color=RGBColor(0xBD, 0xD7, 0xEE))

    # KPI strip at bottom
    strip = s1.shapes.add_shape(1, Inches(0), Inches(5.2), W, Inches(2.3))
    strip.fill.solid(); strip.fill.fore_color.rgb = RGBColor(0x17, 0x2A, 0x4E)
    strip.line.fill.background()

    for i, kpi in enumerate(kpis):
        x = Inches(0.8 + i * 4.18)
        util_color = RED if kpi["util"] > 1 else AMBER if kpi["util"] >= 0.8 else GREEN
        _add_textbox(s1, kpi["label"],
                     x, Inches(5.35), Inches(3.8), Inches(0.45),
                     font_size=14, bold=True, color=WHITE)
        _add_textbox(s1, f"{kpi['dem']:,.0f} samples",
                     x, Inches(5.75), Inches(3.8), Inches(0.6),
                     font_size=22, bold=True, color=util_color)
        _add_textbox(s1, f"{kpi['util']:.0%} of {kpi['cap']:,}/yr",
                     x, Inches(6.35), Inches(3.8), Inches(0.4),
                     font_size=13, color=RGBColor(0xBD, 0xD7, 0xEE))

    # ══════════════════════════════════════════════════════
    #  SLIDE 2 — KPI Summary Table
    # ══════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(blank)
    _solid_bg(s2, LIGHT_GREY)
    _add_title_bar(s2, "KPI Summary — Latest Year Demand & Utilization",
                   f"All figures for {last_year}")

    rows_data = [("Group", "Demand (samples)", "Capacity /yr", "Utilization %", "Status")]
    for kpi in kpis:
        status = "🔴 OVERLOADED" if kpi["util"] > 1 else "🟡 Near Cap" if kpi["util"] >= 0.8 else "🟢 Healthy"
        rows_data.append((
            kpi["label"],
            f"{kpi['dem']:,.0f}",
            f"{kpi['cap']:,}",
            f"{kpi['util']:.1%}",
            status,
        ))

    # Draw table manually as text boxes
    col_widths = [Inches(2.6), Inches(2.2), Inches(2.0), Inches(2.0), Inches(2.2)]
    col_xs     = [Inches(0.4 + sum(col_widths[:i]) / Inches(1) * Inches(1))
                  for i in range(5)]
    # simpler: compute cumulative
    cx = Inches(0.4)
    col_xs = []
    for w in col_widths:
        col_xs.append(cx); cx += w

    row_h  = Inches(0.72)
    row_ys = [Inches(1.3 + i * row_h) for i in range(len(rows_data))]

    hdr_bg_colors = [DARK_BLUE] * 5
    row_bgs = [
        [RGBColor(0xE8, 0xF0, 0xFE)] * 5,  # Mechanical row
        [RGBColor(0xE8, 0xF5, 0xE9)] * 5,  # Coating row
        [RGBColor(0xFF, 0xF3, 0xE0)] * 5,  # Thermal row
    ]

    for ri, row in enumerate(rows_data):
        for ci, (cell_text, col_x, col_w) in enumerate(zip(row, col_xs, col_widths)):
            box = s2.shapes.add_shape(1, col_x, row_ys[ri], col_w - Inches(0.04), row_h - Inches(0.04))
            if ri == 0:
                box.fill.solid(); box.fill.fore_color.rgb = DARK_BLUE
                box.line.fill.background()
                txt_color = WHITE
            else:
                bg = row_bgs[ri - 1][ci]
                box.fill.solid(); box.fill.fore_color.rgb = bg
                box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
                txt_color = DARK_BLUE
                if ci == 4:   # Status column colour-code
                    if "OVER" in cell_text:
                        box.fill.fore_color.rgb = RGBColor(0xFF, 0xCC, 0xCC)
                    elif "Near" in cell_text:
                        box.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xCC)
                    else:
                        box.fill.fore_color.rgb = RGBColor(0xCC, 0xEF, 0xCC)

            _add_textbox(s2, cell_text,
                         col_x + Inches(0.06), row_ys[ri] + Inches(0.1),
                         col_w - Inches(0.12), row_h - Inches(0.15),
                         font_size=13 if ri > 0 else 12,
                         bold=(ri == 0),
                         color=txt_color,
                         align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    # ══════════════════════════════════════════════════════
    #  SLIDE 3 — Comparison Bar Chart
    # ══════════════════════════════════════════════════════
    s3 = prs.slides.add_slide(blank)
    _solid_bg(s3, LIGHT_GREY)
    _add_title_bar(s3, "Annual Demand Comparison — All 3 Lab Groups",
                   "Year-on-year demand vs capacity by group")

    png3 = _fig_to_png(fig_cmp, width=1600, height=700)
    if png3:
        _embed_png(s3, png3, Inches(0.3), Inches(1.2), Inches(12.7), Inches(6.1))
    else:
        _placeholder_box(s3, "Comparison Bar Chart",
                         Inches(0.3), Inches(1.2), Inches(12.7), Inches(6.1))

    # ══════════════════════════════════════════════════════
    #  SLIDE 4 — Utilisation Trend
    # ══════════════════════════════════════════════════════
    s4 = prs.slides.add_slide(blank)
    _solid_bg(s4, LIGHT_GREY)
    _add_title_bar(s4, "Monthly Utilization Trend — All 3 Groups (last 3 years)",
                   "Green <80%  |  Yellow 80–100%  |  Red >100%")

    png4 = _fig_to_png(fig_util, width=1600, height=700)
    if png4:
        _embed_png(s4, png4, Inches(0.3), Inches(1.2), Inches(12.7), Inches(6.1))
    else:
        _placeholder_box(s4, "Utilization Trend",
                         Inches(0.3), Inches(1.2), Inches(12.7), Inches(6.1))

    # ══════════════════════════════════════════════════════
    #  SLIDE 5 — Gantt Charts (3 side by side)
    # ══════════════════════════════════════════════════════
    s5 = prs.slides.add_slide(blank)
    _solid_bg(s5, LIGHT_GREY)
    _add_title_bar(s5, "Gantt — Weekly Lab Occupancy (Current / Latest Year)",
                   "Heatmap: Green=active  Yellow=near cap  Red=overloaded")

    gantt_labels = ["Mechanical", "Coating", "Thermal"]
    gantt_w = Inches(4.2)
    gantt_h = Inches(5.8)
    for gi, (gfig, glabel) in enumerate(zip(figs_gantt, gantt_labels)):
        gx = Inches(0.2 + gi * 4.38)
        gy = Inches(1.15)
        # label
        _add_textbox(s5, f"🔬 {glabel}",
                     gx, gy - Inches(0.3), gantt_w, Inches(0.3),
                     font_size=13, bold=True, color=DARK_BLUE)
        if gfig is not None:
            png5 = _fig_to_png(gfig, width=600, height=800)
            if png5:
                _embed_png(s5, png5, gx, gy, gantt_w, gantt_h)
            else:
                _placeholder_box(s5, f"{glabel} Gantt", gx, gy, gantt_w, gantt_h)
        else:
            _placeholder_box(s5, f"{glabel} Gantt (no data)", gx, gy, gantt_w, gantt_h)

    # ══════════════════════════════════════════════════════
    #  SLIDE 6 — Key Findings (auto-generated)
    # ══════════════════════════════════════════════════════
    s6 = prs.slides.add_slide(blank)
    _solid_bg(s6, LIGHT_GREY)
    _add_title_bar(s6, "Key Findings & Recommendations",
                   f"Auto-computed from {last_year} data")

    # Compute findings
    findings  = []
    risks     = []
    recs      = []

    # Most overloaded group
    max_util_group = max(kpis, key=lambda k: k["util"])
    findings.append(
        f"Most loaded group: {max_util_group['label']} at "
        f"{max_util_group['util']:.0%} utilization "
        f"({max_util_group['dem']:,.0f} / {max_util_group['cap']:,} samples/yr)"
    )

    # Overloaded labs
    all_types_caps = {**caps_a, **caps_b, **caps_c}
    all_wdfs = {}
    for t in types_a: all_wdfs[t] = wdf_a
    for t in types_b: all_wdfs[t] = wdf_b
    for t in types_c: all_wdfs[t] = wdf_c

    all_annual = {}
    all_annual.update(annual_a.get(last_year, {}))
    all_annual.update(annual_b.get(last_year, {}))
    all_annual.update(annual_c.get(last_year, {}))

    for t, dem in all_annual.items():
        cap = all_types_caps.get(t, 1)
        u   = dem / cap if cap > 0 else 0
        if u > 1.0:
            risks.append(f"{t}: {u:.0%} utilization — OVERLOADED ({dem:.0f}/{cap})")
        elif u >= 0.8:
            risks.append(f"{t}: {u:.0%} utilization — Near capacity ({dem:.0f}/{cap})")

    # Highest demand year
    year_totals = {}
    for y in years:
        tot = (sum(annual_a.get(int(y), {}).get(t, 0) for t in types_a) +
               sum(annual_b.get(int(y), {}).get(t, 0) for t in types_b) +
               sum(annual_c.get(int(y), {}).get(t, 0) for t in types_c))
        year_totals[int(y)] = tot
    best_yr = max(year_totals, key=year_totals.get)
    findings.append(
        f"Highest overall demand year: {best_yr} "
        f"({year_totals[best_yr]:,.0f} total samples across all groups)"
    )

    # Recommendations
    for t, dem in all_annual.items():
        cap = all_types_caps.get(t, 1)
        u   = dem / cap if cap > 0 else 0
        if u >= 0.8:
            needed = math.ceil(dem / 0.8)
            recs.append(f"{t}: increase capacity to ≥{needed}/yr (target 80% utilization)")

    # Layout text on slide
    content_y = Inches(1.25)

    def _section(slide, heading, items, y, heading_color=DARK_BLUE, item_color=RGBColor(0x22,0x22,0x22)):
        _add_textbox(slide, heading,
                     Inches(0.45), y, Inches(12.4), Inches(0.45),
                     font_size=16, bold=True, color=heading_color)
        y += Inches(0.42)
        for item in items:
            _add_textbox(slide, f"  •  {item}",
                         Inches(0.55), y, Inches(12.2), Inches(0.38),
                         font_size=13, color=item_color)
            y += Inches(0.36)
        return y + Inches(0.1)

    content_y = _section(s6, "📊 Key Findings", findings, content_y)

    if risks:
        content_y = _section(s6, "⚠️ Risks — Labs at ≥80% Utilization", risks,
                              content_y, heading_color=RGBColor(0xC0,0x00,0x00),
                              item_color=RGBColor(0x88,0x00,0x00))

    if recs:
        _section(s6, "✅ Recommended Capacity Increases", recs, content_y,
                 heading_color=RGBColor(0x1F,0x5C,0x1A),
                 item_color=RGBColor(0x1F,0x5C,0x1A))

    # ── serialize to bytes ────────────────────────────────
    buf = _io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ────────────────────────────────────────────────────────────
#  GUI
# ────────────────────────────────────────────────────────────

